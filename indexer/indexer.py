#!/usr/bin/env python3
"""
TankFinder Indexer — incremental + FTS + QUOTES fallback with newest Q-PDF only

Key features
- JOBS/ARCHIVES: derive job_id from pattern (e.g., 101-23).
- QUOTES: P:\QUOTES\YYYY\...\Q####(.N) => job_id "Q####-YY", q_version=N (0 if none).
- FTS stores text for files; for QUOTES PDFs we keep only the HIGHEST q_version.
- Always roll up job flags (PDF/CAD/COMPRESS/AME) for jobs seen this run.

CLI
  py -3 indexer\indexer.py --year-min 2015 --year-max 2025
  py -3 indexer\indexer.py --quotes-only
  py -3 indexer\indexer.py --rebuild-fts --year-min 2015 --year-max 2025

Config (config.yaml)
  roots: ["P:\\JOBS", "P:\\ARCHIVES\\OLD JOBS"]
  quotes_roots: ["P:\\QUOTES"]
  quotes_scan: {year_min: 2022, year_max: 2100}
  job_id_regex: "(?P<job>\\b\\d{3}-\\d{2}\\b)"
  scan_policy: {...}
  pdf_text: {enabled: true, path_allow_tokens: [...], max_pages: 10, max_chars: 40000}
  office_text: {enabled: false, include: [...], ...}
  ignore: {ext: [...], dir_tokens: [...]}
"""

from __future__ import annotations
import argparse, hashlib, os, re, sqlite3, sys, time
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, Iterator, List, Optional, Tuple
import yaml
from tqdm import tqdm

# -------- Optional PDF text (PyMuPDF) --------
try:
    import fitz  # PyMuPDF
    try:
        fitz.TOOLS.mupdf_display_errors(False)  # quiet noisy warnings
    except Exception:
        pass
except Exception:
    fitz = None

# -------- Optional Office parsing (off by default via YAML) --------
try:
    import openpyxl
except Exception:
    openpyxl = None
try:
    from docx import Document
except Exception:
    Document = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None

# -------- Paths / config helpers --------
def project_root() -> Path:
    # indexer is run as a script; its project root is one up from /indexer
    return Path(__file__).resolve().parents[1]

ROOT        = project_root()
CFG_PATH    = ROOT / "config.yaml"
SCHEMA_PATH = Path(__file__).with_name("schema.sql")
DEFAULT_DB  = ROOT / "tankfinder.db"

def load_cfg() -> dict:
    with open(CFG_PATH, "r", encoding="utf-8") as f:
        return yaml.safe_load(f) or {}

def resolve_db_path(cfg: dict) -> Path:
    """
    Order of precedence:
    1) TANKFINDER_DB env var (what the GUI uses in production)
    2) cfg['db_path'] if provided
    3) default next to the repo (DEFAULT_DB)
    """
    env = os.getenv("TANKFINDER_DB")
    if env:
        return Path(env)
    cfg_db = (cfg or {}).get("db_path")
    if cfg_db:
        return Path(cfg_db)
    return DEFAULT_DB

# -------- Quote / Year patterns --------
# e.g., Q9185, Q9185.2 — with word boundaries so we don't match part of another token
QNUM_RE     = re.compile(r"(?i)(?<!\w)q(?P<num>\d{4,6})(?:\.(?P<ver>\d+))?(?!\w)")
YEAR_DIR_RE = re.compile(r"^(19|20)\d{2}$", re.I)


# ================= util =================
def utc_iso(ts: float) -> str:
    return datetime.fromtimestamp(ts, tz=timezone.utc).isoformat()

def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()

def file_hash16(s: str) -> str:
    return hashlib.sha1(s.encode("utf-8", errors="ignore")).hexdigest()[:16]

_NON_ALNUM = re.compile(r"[^a-z0-9]+")
def norm_tokens(s: str) -> List[str]:
    return [t for t in _NON_ALNUM.split(s.lower()) if t]

def load_cfg() -> dict:
    with open(CFG_PATH, "r", encoding="utf-8") as f:
        return yaml.safe_load(f) or {}

def connect_db(db_path: Path) -> sqlite3.Connection:
    con = sqlite3.connect(db_path)
    con.execute("PRAGMA journal_mode=WAL;")
    con.execute("PRAGMA synchronous=NORMAL;")
    con.execute("PRAGMA temp_store=MEMORY;")
    con.execute("PRAGMA foreign_keys=ON;")
    return con


def ensure_schema(con: sqlite3.Connection, rebuild_fts: bool = False) -> None:
    # Apply schema.sql if present
    if SCHEMA_PATH.exists():
        with open(SCHEMA_PATH, "r", encoding="utf-8") as f:
            con.executescript(f.read())

    # Add missing columns / indexes we rely on
    jobs_cols = {r[1] for r in con.execute("PRAGMA table_info(jobs)")}
    if "job_year" not in jobs_cols:
        try: con.execute("ALTER TABLE jobs ADD COLUMN job_year INTEGER;")
        except sqlite3.OperationalError: pass
    con.execute("CREATE INDEX IF NOT EXISTS idx_jobs_year ON jobs(job_year);")
    con.execute("CREATE INDEX IF NOT EXISTS idx_jobs_flags ON jobs(has_compress, has_ame, has_dwg_dxf, has_pdf);")

    files_cols = {r[1] for r in con.execute("PRAGMA table_info(files)")}
    if "q_version" not in files_cols:
        try: con.execute("ALTER TABLE files ADD COLUMN q_version INTEGER;")
        except sqlite3.OperationalError: pass
    con.execute("CREATE INDEX IF NOT EXISTS idx_files_job_del ON files(job_id, deleted);")
    con.execute("CREATE INDEX IF NOT EXISTS idx_files_job_ext_del ON files(job_id, ext, deleted);")
    con.execute("CREATE INDEX IF NOT EXISTS idx_files_hash16 ON files(file_hash16);")
    con.execute("CREATE INDEX IF NOT EXISTS idx_files_job_qver ON files(job_id, q_version);")

    # FTS
    has_fts = bool(con.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='fts_files'").fetchone())
    if rebuild_fts and has_fts:
        con.execute("DROP TABLE fts_files;")
        has_fts = False
    if not has_fts:
        con.execute("""
            CREATE VIRTUAL TABLE fts_files USING fts5(
              content,
              file_hash16 UNINDEXED,
              tokenize = "unicode61 separators '-_./()[]{}' remove_diacritics 2"
            );
        """)
    con.commit()

# ================= detectors =================
DEFAULT_DETECTORS = {
    "compress":   {"ext_any": {".cw7", ".xml"},                  "name_tokens_any": {"compress","codeware"}},
    "ametank":    {"ext_any": {".mdl", ".xmt_txt"},              "name_tokens_any": {"ametank","ame"}},
    "cad":        {"ext_any": {".dwg", ".dxf"}},
    "pdf":        {"ext_any": {".pdf"}},
    "excel":      {"ext_any": {".xlsx",".xlsm",".xls",".csv"}},
    "word":       {"ext_any": {".docx",".doc"}},
    "powerpoint": {"ext_any": {".pptx",".ppt"}},
    "archive":    {"ext_any": {".zip",".7z",".rar"}},
    "legacy_calc":{"ext_any": {".wk1",".wk3",".wk4",".fmt",".prn"}},
}
def load_detectors(cfg: dict) -> Dict[str, dict]:
    det = {k: {kk: set(vv) if isinstance(vv,(list,set,tuple)) else vv for kk,vv in v.items()} for k,v in DEFAULT_DETECTORS.items()}
    for key, spec in (cfg.get("detectors") or {}).items():
        d = det.setdefault(key, {})
        if "ext_any" in spec: d["ext_any"] = set(spec["ext_any"])
        if "name_tokens_any" in spec: d["name_tokens_any"] = set(t.lower() for t in spec["name_tokens_any"])
    return det

# ================= data model =================
@dataclass
class FileRow:
    file_hash16: str
    job_id: str
    rel_path: str
    ext: str
    size_bytes: int
    mtime_utc: str
    kind: str
    tokens_fname: str
    detector_hits: str
    deleted: int = 0
    q_version: Optional[int] = None

JOB_ID_PAT: Optional[re.Pattern] = None

def parse_job_id_from_path(path: Path, job_re: re.Pattern) -> Optional[str]:
    for p in [path, *path.parents]:
        m = job_re.search(str(p))
        if m: return m.group("job")
    return None

def job_year_from_job_id(job_id: str) -> Optional[int]:
    try:
        yy = int(job_id.split("-")[-1])
        return 1900 + yy if yy >= 90 else 2000 + yy
    except Exception:
        return None

def detect_kind(ext: str) -> str:
    e = ext.lower()
    if e == ".pdf": return "pdf"
    if e in {".dwg",".dxf"}: return "cad"
    if e in {".jpg",".jpeg",".png",".bmp",".tif",".tiff",".heic"}: return "image"
    if e in {".txt",".csv",".log",".md",".xml",".html",".htm"}: return "text"
    return "other"

def apply_detectors(tokens: List[str], ext: str, detectors: Dict[str,dict]) -> List[str]:
    hits = []; tokset = set(tokens); e = ext.lower()
    for label, spec in detectors.items():
        ext_any = spec.get("ext_any") or set()
        nt_any  = spec.get("name_tokens_any") or set()
        if (ext_any and e in ext_any) or (nt_any and (tokset & nt_any)):
            hits.append(label)
    return hits

# ================= QUOTES helpers =================
def parse_quote_context(path: Path, quotes_roots: List[str], year_min: int, year_max: int):
    r"""
    If path is under QUOTES\YYYY\... and name or a parent contains Q####(.N),
    return (job_id, job_year, job_root, q_version). Else None.
    """
    try:
        rp = path.resolve()
    except Exception:
        rp = path

    for root in quotes_roots:
        try:
            rel = rp.relative_to(Path(root).resolve())
        except Exception:
            continue

        parts = rel.parts
        if not parts:
            continue
        y = parts[0]
        if not YEAR_DIR_RE.fullmatch(y):
            continue
        yint = int(y)
        if not (year_min <= yint <= year_max):
            continue

        # Q####(.N) in parent chain or filename
        q_folder = None
        for pr in [path.parent, *path.parents]:
            if QNUM_RE.search(pr.name):
                q_folder = pr
                break
        m = QNUM_RE.search((q_folder.name if q_folder else path.stem))
        if not m:
            continue

        num = m.group("num")
        ver = m.group("ver")
        qver = int(ver) if ver else 0

        job_id  = f"Q{num}-{yint % 100:02d}"
        job_root = q_folder or path.parent
        return job_id, yint, job_root, qver

    return None

def should_parse_pdf_jobs(path: Path, cfg: dict) -> bool:
    pdf_cfg = (cfg.get("pdf_text") or {})
    if not pdf_cfg.get("enabled") or fitz is None: return False
    if path.suffix.lower() != ".pdf": return False
    allow = set(x.lower() for x in pdf_cfg.get("path_allow_tokens", []))
    if not allow: return True
    return any(tok in str(path.parent).lower() for tok in allow)

def should_parse_pdf_quotes_only(path: Path, quotes_roots: List[str], year_min: int, year_max: int) -> bool:
    if path.suffix.lower() != ".pdf": return False
    # Only PDFs with Q####(.N) in folder or filename and under a valid year
    qc = parse_quote_context(path, quotes_roots, year_min, year_max)
    return qc is not None

def extract_pdf_text(path: Path, max_pages=10, max_chars=40000) -> str:
    if fitz is None: return ""
    try:
        doc = fitz.open(str(path))
        pages = min(max_pages, len(doc))
        chunks = []
        for i in range(pages):
            chunks.append(doc[i].get_text("text"))
        doc.close()
        txt = " ".join(" ".join(chunks).split())
        return txt[:max_chars]
    except Exception:
        return ""

# ---- office helpers (optional) ----
def read_text_file_safe(path: Path, max_chars: int = 40000) -> str:
    for enc in ("utf-8","latin-1"):
        try:
            with open(path, "r", encoding=enc, errors="ignore") as f:
                return f.read(max_chars)
        except Exception:
            continue
    return ""

def extract_csv_text(path: Path, max_lines: int = 200, max_chars: int = 40000) -> str:
    try:
        out = []
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            for i,line in enumerate(f):
                if i >= max_lines: break
                out.append(line.strip())
        return " ".join(out)[:max_chars]
    except Exception:
        return ""

def extract_xlsx_text(path: Path, sheet_limit=3, cell_limit=500, max_chars=40000) -> str:
    if openpyxl is None: return ""
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        chunks = []
        for si, ws in enumerate(wb.worksheets):
            if si >= sheet_limit: break
            count = 0
            for row in ws.iter_rows(values_only=True):
                for v in row:
                    if v is None: continue
                    chunks.append(str(v)); count += 1
                    if count >= cell_limit: break
                if count >= cell_limit: break
        wb.close()
        return " ".join(chunks)[:max_chars]
    except Exception:
        return ""

def extract_docx_text(path: Path, max_chars=40000) -> str:
    if Document is None: return ""
    try:
        doc = Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text]
        for t in doc.tables:
            for row in t.rows:
                for cell in row.cells:
                    if cell.text: parts.append(cell.text)
        return " ".join(" ".join(parts).split())[:max_chars]
    except Exception:
        return ""

def extract_pptx_text(path: Path, max_chars=40000) -> str:
    if Presentation is None: return ""
    try:
        prs = Presentation(str(path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape,"text") and shape.text:
                    parts.append(shape.text)
        return " ".join(" ".join(parts).split())[:max_chars]
    except Exception:
        return ""

def extract_office_text(path: Path, cfg: dict) -> str:
    oc = (cfg.get("office_text") or {})
    if not oc.get("enabled"): return ""
    include = set(x.lower() for x in oc.get("include", []))
    max_chars = int(oc.get("max_chars", 40000))
    ext = path.suffix.lower()
    if ext == ".csv" and "csv" in include:
        return extract_csv_text(path, max_lines=int(oc.get("csv_max_lines",200)), max_chars=max_chars)
    if ext in (".xlsx",".xlsm") and "xlsx" in include:
        return extract_xlsx_text(path, sheet_limit=int(oc.get("xlsx_sheet_limit",3)),
                                 cell_limit=int(oc.get("xlsx_cells_limit",500)),
                                 max_chars=max_chars)
    if ext == ".docx" and "docx" in include:
        return extract_docx_text(path, max_chars=max_chars)
    if ext == ".pptx" and "pptx" in include:
        return extract_pptx_text(path, max_chars=max_chars)
    if ext in (".txt",".md",".log",".xml",".html",".htm"):
        return read_text_file_safe(path, max_chars=max_chars)
    return ""

# ================= scanning =================
def walk_files(roots: List[str], scan_policy: dict, denylist_paths: List[str]) -> Iterator[Path]:
    deny = [p.lower().rstrip("\\/") for p in (denylist_paths or [])]
    year_only = bool((scan_policy or {}).get("only_year_dirs_under_roots", False))
    year_re   = re.compile((scan_policy or {}).get("year_dir_regex", r"^\d{4}$"), re.I)
    year_min  = int((scan_policy or {}).get("year_min", 1900))
    year_max  = int((scan_policy or {}).get("year_max", 2100))

    def denied(path: Path) -> bool:
        plow = str(path).lower(); return any(plow.startswith(p) for p in deny)

    def push_children_year_dirs(rootp: Path, stack: List[Path]) -> None:
        try:
            with os.scandir(rootp) as it:
                for e in it:
                    if not e.is_dir(follow_symlinks=False): continue
                    name = e.name
                    if not year_re.match(name): continue
                    try: yr = int(re.findall(r"\d{4}", name)[0])
                    except Exception: continue
                    if yr < year_min or yr > year_max: continue
                    p = Path(e.path)
                    if not denied(p): stack.append(p)
        except (PermissionError, FileNotFoundError, OSError):
            return

    def scandir_safe(d: Path):
        for attempt in (0,1):
            try: return os.scandir(d)
            except (FileNotFoundError, OSError):
                if attempt==0: time.sleep(0.05); continue
                return None
            except PermissionError: return None

    for root in roots:
        rootp = Path(root)
        if not rootp.exists() or denied(rootp): continue
        stack: List[Path] = []
        if year_only: push_children_year_dirs(rootp, stack)
        else: stack.append(rootp)
        while stack:
            d = stack.pop()
            it = scandir_safe(d)
            if it is None: continue
            try:
                with it:
                    for e in it:
                        try:
                            if e.is_dir(follow_symlinks=False):
                                name = e.name.lower()
                                if name in {"$recycle.bin","system volume information"}: continue
                                p = Path(e.path)
                                if not denied(p): stack.append(p)
                            elif e.is_file(follow_symlinks=False):
                                yield Path(e.path)
                        except (PermissionError, FileNotFoundError, OSError):
                            continue
            except (PermissionError, FileNotFoundError, OSError):
                continue

# ================= db ops =================
def ensure_job(con: sqlite3.Connection, job_id: str, root_path: str, job_year: Optional[int]) -> None:
    now = now_iso()
    con.execute("""
      INSERT INTO jobs (job_id, root_path, job_year, first_seen, last_seen)
      VALUES (?, ?, ?, ?, ?)
      ON CONFLICT(job_id) DO UPDATE SET
        root_path=excluded.root_path,
        job_year=COALESCE(jobs.job_year, excluded.job_year),
        last_seen=excluded.last_seen
    """, (job_id, root_path, job_year, now, now))

def upsert_files(con: sqlite3.Connection, rows: List[FileRow]) -> None:
    con.executemany("""
      INSERT INTO files (file_hash16, job_id, rel_path, ext, size_bytes, mtime_utc,
                         kind, tokens_fname, detector_hits, deleted, q_version)
      VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
      ON CONFLICT(file_hash16) DO UPDATE SET
        size_bytes=excluded.size_bytes,
        mtime_utc=excluded.mtime_utc,
        kind=excluded.kind,
        tokens_fname=excluded.tokens_fname,
        detector_hits=excluded.detector_hits,
        deleted=0,
        q_version=COALESCE(excluded.q_version, files.q_version)
    """, [(r.file_hash16, r.job_id, r.rel_path, r.ext, r.size_bytes, r.mtime_utc,
           r.kind, r.tokens_fname, r.detector_hits, r.deleted, r.q_version)
          for r in rows])

def upsert_fts_rows(con: sqlite3.Connection, fts_rows: List[Tuple[str,str]]) -> None:
    if not fts_rows: return
    for content, fh in fts_rows:
        con.execute("DELETE FROM fts_files WHERE file_hash16=?", (fh,))
        con.execute("INSERT INTO fts_files(content, file_hash16) VALUES (?,?)", (content, fh))

def mark_deleted_missing(con: sqlite3.Connection, seen_hashes: set, year_min: Optional[int], year_max: Optional[int]) -> int:
    if year_min is not None and year_max is not None:
        cur = con.execute("""
          SELECT f.file_hash16
          FROM files f JOIN jobs j ON j.job_id=f.job_id
          WHERE f.deleted=0 AND j.job_year BETWEEN ? AND ?
        """, (year_min, year_max))
    else:
        cur = con.execute("SELECT file_hash16 FROM files WHERE deleted=0")
    to_delete = [(fh,) for (fh,) in cur.fetchall() if fh not in seen_hashes]
    if to_delete:
        con.executemany("UPDATE files SET deleted=1 WHERE file_hash16=?", to_delete)
    return len(to_delete)

def rollup_job_stats(con: sqlite3.Connection, job_id: str) -> None:
    fc, bytes_, maxmt = con.execute(
        "SELECT COUNT(*), COALESCE(SUM(size_bytes),0), MAX(mtime_utc) FROM files WHERE job_id=? AND deleted=0",
        (job_id,)
    ).fetchone()
    def exists(sql): return con.execute(sql, (job_id,)).fetchone()[0]
    has_pdf      = exists("SELECT EXISTS(SELECT 1 FROM files WHERE job_id=? AND deleted=0 AND (ext='.pdf' OR instr(detector_hits,'pdf')>0))")
    has_cad      = exists("SELECT EXISTS(SELECT 1 FROM files WHERE job_id=? AND deleted=0 AND (ext IN('.dwg','.dxf') OR instr(detector_hits,'cad')>0))")
    has_compress = exists("SELECT EXISTS(SELECT 1 FROM files WHERE job_id=? AND deleted=0 AND instr(detector_hits,'compress')>0)")
    has_ame      = exists("SELECT EXISTS(SELECT 1 FROM files WHERE job_id=? AND deleted=0 AND instr(detector_hits,'ametank')>0)")
    try:
        has_legacy = exists("SELECT EXISTS(SELECT 1 FROM files WHERE job_id=? AND deleted=0 AND instr(detector_hits,'legacy_calc')>0)")
        con.execute("""UPDATE jobs SET
          file_count_total=?, byte_size_total=?, has_pdf=?, has_dwg_dxf=?, has_compress=?, has_ame=?, has_legacy_calc=?, last_modified_utc=?
          WHERE job_id=?""", (fc, bytes_, has_pdf, has_cad, has_compress, has_ame, has_legacy, maxmt, job_id))
    except sqlite3.OperationalError:
        con.execute("""UPDATE jobs SET
          file_count_total=?, byte_size_total=?, has_pdf=?, has_dwg_dxf=?, has_compress=?, has_ame=?, last_modified_utc=?
          WHERE job_id=?""", (fc, bytes_, has_pdf, has_cad, has_compress, has_ame, maxmt, job_id))

def cleanup_old_quote_versions(con: sqlite3.Connection):
    # Keep only newest Q-PDF version per quote in FTS
    con.execute("""
      DELETE FROM fts_files
      WHERE file_hash16 IN (
        SELECT f.file_hash16
        FROM files f
        JOIN (
          SELECT job_id, MAX(COALESCE(q_version,0)) AS maxv
          FROM files
          WHERE job_id LIKE 'Q%' AND ext='.pdf' AND deleted=0
          GROUP BY job_id
        ) m ON f.job_id=m.job_id
        WHERE f.job_id LIKE 'Q%' AND f.ext='.pdf' AND f.deleted=0
          AND COALESCE(f.q_version,0) < m.maxv
      );
    """)
    con.commit()

# ================= main =================
def main():
    ap = argparse.ArgumentParser(description="TankFinder Indexer")
    ap.add_argument("--limit", type=int, default=0, help="Stop after N files (debug/partial)")
    ap.add_argument("--dry-run", action="store_true", help="Scan but do not write DB")
    ap.add_argument("--no-delete", action="store_true", help="Skip delete pass (safety)")
    ap.add_argument("--year-min", type=int, default=2015, help="Min JOB year to include (via job_id)")
    ap.add_argument("--year-max", type=int, default=2025, help="Max JOB year to include (via job_id)")
    ap.add_argument("--rebuild-fts", action="store_true", help="Drop and recreate FTS table")
    ap.add_argument("--quotes-only", action="store_true", help="Scan only quotes_roots (skip JOBS/ARCHIVES)")
    args = ap.parse_args()

    cfg = load_cfg()

    # --- DB path (env > YAML > default) ---
    db_path = resolve_db_path(cfg)          # <-- NEW
    print(f"[indexer] DB -> {db_path}")     # helpful when running from a share
    db_path.parent.mkdir(parents=True, exist_ok=True)

    # ignore config
    ignore_cfg  = (cfg.get("ignore") or {})
    ignore_exts = {e.lower() for e in ignore_cfg.get("ext", [])}
    ignore_dir_tokens = {t.lower() for t in ignore_cfg.get("dir_tokens", [])}

    # job id regex
    global JOB_ID_PAT
    JOB_ID_PAT = re.compile(cfg.get("job_id_regex") or r"(?P<job>\b\d{3}-\d{2}\b)", re.I)

    detectors     = load_detectors(cfg)

    roots         = cfg.get("roots") or []
    quotes_roots  = cfg.get("quotes_roots") or []

    if args.quotes_only:
        scan_roots = quotes_roots
        args.no_delete = True  # never delete on subset scans
    else:
        scan_roots = roots + quotes_roots

    if not scan_roots:
        print("No roots configured in config.yaml", file=sys.stderr); sys.exit(2)

    scan_policy   = (cfg.get("scan_policy") or {})
    denylist      = (cfg.get("denylist_paths") or {})

    qs_cfg        = (cfg.get("quotes_scan") or {})
    q_year_min    = int(qs_cfg.get("year_min", 2022))
    q_year_max    = int(qs_cfg.get("year_max", 2100))

    # --- open DB using the resolved path ---
    con = connect_db(db_path)               # <-- CHANGED
    ensure_schema(con, rebuild_fts=args.rebuild_fts)

    pdf_cfg = (cfg.get("pdf_text") or {})
    max_pdf_pages = int(pdf_cfg.get("max_pages", 10))
    max_pdf_chars = int(pdf_cfg.get("max_chars", 40000))

    seen_hashes: set[str] = set()
    batch: List[FileRow]  = []
    fts_batch: List[Tuple[str, str]] = []
    per_job_seen_roots: Dict[str, str] = {}
    counters = {
        "total_scanned": 0, "indexed": 0, "fts_backfilled": 0,
        "skipped_no_job": 0, "skipped_unchanged": 0, "skipped_out_of_year": 0,
        "skipped_ignored_ext": 0, "skipped_ignored_dir": 0
    }


    start = time.time()
    for p in tqdm(walk_files(scan_roots, scan_policy, denylist), desc="Scanning"):
        counters["total_scanned"] += 1
        if args.limit and counters["total_scanned"] > args.limit:
            break

        # quick cheap skips
        if p.suffix.lower() in ignore_exts:
            counters["skipped_ignored_ext"] += 1; continue
        if ignore_dir_tokens and any(tok in str(p.parent).lower() for tok in ignore_dir_tokens):
            counters["skipped_ignored_dir"] += 1; continue

        try:
            st = p.stat()
        except (FileNotFoundError, PermissionError, OSError):
            continue
        size = st.st_size
        ext  = p.suffix.lower()

        # ---------- resolve job_id / job_root / jy (with QUOTES fallback) ----------
        job_id   = parse_job_id_from_path(p, JOB_ID_PAT)
        job_root = None
        jy       = None
        qver     = None

        if not job_id:
            qc = parse_quote_context(p, [str(r) for r in quotes_roots], q_year_min, q_year_max)
            if qc:
                job_id, jy, job_root, qver = qc

        if not job_id:
            counters["skipped_no_job"] += 1
            continue

        # JOBS year gate (quotes use their own window)
        if job_id and not job_id.startswith("Q"):
            jy = jy if jy is not None else job_year_from_job_id(job_id)
            if jy is not None and ((jy < args.year_min) or (jy > args.year_max)):
                counters["skipped_out_of_year"] += 1; continue

        # Find job_root if it's a JOBS item
        if job_root is None:
            for parent in [p.parent, *p.parents]:
                if JOB_ID_PAT.search(str(parent)):
                    job_root = parent; break
        if job_root is None:
            # Fallback root for QUOTES already set; or skip if cannot determine
            job_root = p.parent

        rel = str(p).replace(str(job_root) + os.sep, "", 1) if str(p).startswith(str(job_root)+os.sep) else p.name
        fh  = file_hash16(str(p).lower())
        mtime_iso = utc_iso(st.st_mtime)

        # unchanged fast-path (+ FTS backfill if needed)
        row = con.execute("SELECT size_bytes, mtime_utc FROM files WHERE file_hash16=?", (fh,)).fetchone()
        if row:
            old_size, old_mtime = row
            if int(old_size) == int(size) and old_mtime == mtime_iso:
                if not args.dry_run:
                    has_fts = con.execute("SELECT 1 FROM fts_files WHERE file_hash16=? LIMIT 1", (fh,)).fetchone()
                    if not has_fts:
                        name_tokens = norm_tokens(p.name) + norm_tokens(str(p.parent))
                        fts_content = " ".join(name_tokens[:64])
                        parse_jobs_pdf   = should_parse_pdf_jobs(p, cfg) if not job_id.startswith("Q") else False
                        parse_quotes_pdf = should_parse_pdf_quotes_only(p, [str(r) for r in quotes_roots], q_year_min, q_year_max) if job_id.startswith("Q") else False
                        if parse_jobs_pdf or parse_quotes_pdf:
                            txt = extract_pdf_text(p, max_pages=max_pdf_pages, max_chars=max_pdf_chars)
                            if txt: fts_content = (fts_content + " " + txt).strip()
                        office_txt = extract_office_text(p, cfg)
                        if office_txt: fts_content = (fts_content + " " + office_txt).strip()
                        fts_batch.append((fts_content, fh)); counters["fts_backfilled"] += 1
                        if len(fts_batch) >= 800:
                            upsert_fts_rows(con, fts_batch); con.commit(); fts_batch.clear()
                seen_hashes.add(fh); counters["skipped_unchanged"] += 1
                continue

        # ensure job row
        if job_id not in per_job_seen_roots:
            per_job_seen_roots[job_id] = str(job_root)
            if not args.dry_run:
                ensure_job(con, job_id, str(job_root), jy)

        # insert/queue this file
        name_tokens = norm_tokens(p.name) + norm_tokens(str(p.parent))
        hits = apply_detectors(name_tokens, ext, detectors)
        kind = detect_kind(ext)

        fr = FileRow(
            file_hash16=fh, job_id=job_id, rel_path=rel, ext=ext, size_bytes=size,
            mtime_utc=mtime_iso, kind=kind, tokens_fname=" ".join(name_tokens[:64]),
            detector_hits=",".join(hits), q_version=(qver if (job_id.startswith("Q") and ext==".pdf") else None)
        )
        batch.append(fr); seen_hashes.add(fh); counters["indexed"] += 1

        # FTS content
        fts_content = fr.tokens_fname
        parse_jobs_pdf   = should_parse_pdf_jobs(p, cfg) if not job_id.startswith("Q") else False
        parse_quotes_pdf = should_parse_pdf_quotes_only(p, [str(r) for r in quotes_roots], q_year_min, q_year_max) if job_id.startswith("Q") else False
        if parse_jobs_pdf or parse_quotes_pdf:
            txtc = extract_pdf_text(p, max_pages=max_pdf_pages, max_chars=max_pdf_chars)
            if txtc: fts_content = (fts_content + " " + txtc).strip()
        office_txt = extract_office_text(p, cfg)
        if office_txt: fts_content = (fts_content + " " + office_txt).strip()
        if not args.dry_run:
            fts_batch.append((fts_content, fh))

        if len(batch) >= 800 and not args.dry_run:
            upsert_files(con, batch); upsert_fts_rows(con, fts_batch); con.commit()
            batch.clear(); fts_batch.clear()

    # tail flush
    if batch and not args.dry_run:
        upsert_files(con, batch); upsert_fts_rows(con, fts_batch); con.commit()
        batch.clear(); fts_batch.clear()

    # delete-pass only on complete scans of all roots (safety)
    deleted = 0
    complete_scan = (args.limit == 0)
    if not args.dry_run and complete_scan and not args.no_delete and not args.quotes_only:
        deleted = mark_deleted_missing(con, seen_hashes, args.year_min, args.year_max)

    # ALWAYS roll up flags for jobs touched this run (includes quotes-only)
    if not args.dry_run:
        for job_id in per_job_seen_roots.keys():
            rollup_job_stats(con, job_id)

        # Enforce newest Q-PDF only in FTS
        cleanup_old_quote_versions(con)

    dur = time.time() - start
    print(f"[OK] Scanned {counters['total_scanned']:,}; indexed={counters['indexed']:,}; "
          f"fts_backfilled={counters['fts_backfilled']:,}; skipped_no_job={counters['skipped_no_job']:,}; "
          f"skipped_unchanged={counters['skipped_unchanged']:,}; skipped_out_of_year={counters['skipped_out_of_year']:,}; "
          f"skipped_ignored_ext={counters['skipped_ignored_ext']:,}; skipped_ignored_dir={counters['skipped_ignored_dir']:,}; "
          f"deleted_marked={deleted:,} in {dur:,.1f}s")

if __name__ == "__main__":
    main()
